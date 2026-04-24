"""
AI Resume Rewriter — Upload PDF/DOCX/PPTX → AI rewrites → Download updated DOCX
"""
import streamlit as st
import json, io, os, tempfile
from core.store import get_settings, get_prefs


# ── Text extractors ───────────────────────────────────────────

def extract_pdf(file_bytes: bytes) -> str:
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def extract_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = filename.lower().split(".")[-1]
    if ext == "pdf":
        return extract_pdf(file_bytes)
    elif ext in ("docx", "doc"):
        return extract_docx(file_bytes)
    elif ext in ("pptx", "ppt"):
        return extract_pptx(file_bytes)
    else:
        try:
            return file_bytes.decode("utf-8", errors="ignore")
        except Exception:
            return ""


# ── AI rewrite via Groq ───────────────────────────────────────

def ai_rewrite_resume(resume_text: str, jd_text: str, name: str,
                       groq_key: str, mode: str = "full") -> dict:
    """
    mode: 'full' = rewrite entire resume  |  'targeted' = targeted sections only
    Returns dict with sections: summary, experience, skills, education, extra_sections
    """
    prompt = f"""You are an expert resume writer and ATS specialist.

ORIGINAL RESUME:
{resume_text[:4000]}

{"JOB DESCRIPTION TO TARGET:" if jd_text else ""}
{jd_text[:1500] if jd_text else "Optimize for general job market."}

CANDIDATE NAME: {name}

{"Rewrite the ENTIRE resume" if mode == "full" else "Rewrite only the Summary and Skills sections"} to:
1. Use strong action verbs and quantified achievements
2. Optimize for ATS keyword matching
3. Keep all real experience — never invent facts
4. Make it concise, professional, and impactful
{"5. Tailor specifically to the job description above" if jd_text else ""}

Return ONLY valid JSON (no markdown, no backticks):
{{
  "candidate_name": "...",
  "professional_summary": "3-4 sentence punchy summary",
  "work_experience": [
    {{
      "title": "Job Title",
      "company": "Company Name",
      "duration": "Month Year – Month Year",
      "bullets": ["Achievement 1 with metric", "Achievement 2", "Achievement 3"]
    }}
  ],
  "skills": ["skill1", "skill2", "skill3", ...],
  "education": [
    {{
      "degree": "Degree Name",
      "institution": "University",
      "year": "Year"
    }}
  ],
  "certifications": ["cert1", "cert2"],
  "keywords_added": ["new keyword 1", "new keyword 2"]
}}"""

    import urllib.request
    payload = json.dumps({
        "model": "llama-3.1-8b-instant",
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": 2000,
        "temperature": 0.3,
    }).encode()
    req = urllib.request.Request(
        "https://api.groq.com/openai/v1/chat/completions",
        data=payload,
        headers={"Authorization": f"Bearer {groq_key}", "Content-Type": "application/json"},
        method="POST"
    )
    with urllib.request.urlopen(req, timeout=40) as r:
        raw = json.loads(r.read())["choices"][0]["message"]["content"]
    raw = raw.strip().strip("```json").strip("```").strip()
    return json.loads(raw)


# ── Build DOCX from AI output ─────────────────────────────────

def build_docx(data: dict) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)

    def add_horizontal_rule(doc):
        para = doc.add_paragraph()
        pPr = para._p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        bottom = OxmlElement('w:bottom')
        bottom.set(qn('w:val'), 'single')
        bottom.set(qn('w:sz'), '6')
        bottom.set(qn('w:space'), '1')
        bottom.set(qn('w:color'), '22d3a5')
        pBdr.append(bottom)
        pPr.append(pBdr)
        para.paragraph_format.space_after = Pt(2)
        return para

    def set_font(run, size=11, bold=False, color=None):
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            r, g, b = int(color[0:2],16), int(color[2:4],16), int(color[4:6],16)
            run.font.color.rgb = RGBColor(r, g, b)

    # ── Header: Name ──
    name_para = doc.add_paragraph()
    name_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    name_run = name_para.add_run(data.get("candidate_name", "Resume"))
    set_font(name_run, size=22, bold=True, color="0a0e1a")
    name_para.paragraph_format.space_after = Pt(4)

    # ── Professional Summary ──
    add_horizontal_rule(doc)
    sec = doc.add_paragraph()
    sec_run = sec.add_run("PROFESSIONAL SUMMARY")
    set_font(sec_run, size=10, bold=True, color="22d3a5")
    sec.paragraph_format.space_after = Pt(3)

    summary_para = doc.add_paragraph()
    s_run = summary_para.add_run(data.get("professional_summary", ""))
    set_font(s_run, size=11)
    summary_para.paragraph_format.space_after = Pt(8)

    # ── Work Experience ──
    add_horizontal_rule(doc)
    sec = doc.add_paragraph()
    sec_run = sec.add_run("WORK EXPERIENCE")
    set_font(sec_run, size=10, bold=True, color="22d3a5")
    sec.paragraph_format.space_after = Pt(3)

    for exp in data.get("work_experience", []):
        title_para = doc.add_paragraph()
        title_run = title_para.add_run(exp.get("title", ""))
        set_font(title_run, size=12, bold=True)
        title_para.paragraph_format.space_after = Pt(1)

        company_para = doc.add_paragraph()
        co_run = company_para.add_run(f"{exp.get('company', '')}  |  {exp.get('duration', '')}")
        set_font(co_run, size=10, color="5a7090")
        company_para.paragraph_format.space_after = Pt(3)

        for bullet in exp.get("bullets", []):
            b_para = doc.add_paragraph(style="List Bullet")
            b_run = b_para.add_run(bullet)
            set_font(b_run, size=11)
            b_para.paragraph_format.space_after = Pt(2)

        doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # ── Skills ──
    add_horizontal_rule(doc)
    sec = doc.add_paragraph()
    sec_run = sec.add_run("SKILLS")
    set_font(sec_run, size=10, bold=True, color="22d3a5")
    sec.paragraph_format.space_after = Pt(3)

    skills = data.get("skills", [])
    if skills:
        skills_para = doc.add_paragraph()
        skills_run = skills_para.add_run("  ·  ".join(skills))
        set_font(skills_run, size=11)
        skills_para.paragraph_format.space_after = Pt(8)

    # ── Education ──
    edu_list = data.get("education", [])
    if edu_list:
        add_horizontal_rule(doc)
        sec = doc.add_paragraph()
        sec_run = sec.add_run("EDUCATION")
        set_font(sec_run, size=10, bold=True, color="22d3a5")
        sec.paragraph_format.space_after = Pt(3)

        for edu in edu_list:
            edu_para = doc.add_paragraph()
            edu_run = edu_para.add_run(
                f"{edu.get('degree', '')}  —  {edu.get('institution', '')}  |  {edu.get('year', '')}"
            )
            set_font(edu_run, size=11)
            edu_para.paragraph_format.space_after = Pt(3)

    # ── Certifications ──
    certs = data.get("certifications", [])
    if certs:
        add_horizontal_rule(doc)
        sec = doc.add_paragraph()
        sec_run = sec.add_run("CERTIFICATIONS")
        set_font(sec_run, size=10, bold=True, color="22d3a5")
        sec.paragraph_format.space_after = Pt(3)
        for cert in certs:
            c_para = doc.add_paragraph()
            c_run = c_para.add_run(f"• {cert}")
            set_font(c_run, size=11)
            c_para.paragraph_format.space_after = Pt(2)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


# ── Main render ───────────────────────────────────────────────

def render():
    settings = get_settings()
    prefs = get_prefs()
    try:
        groq_key = settings.get("groq_api_key", "") or st.secrets.get("GROQ_API_KEY", "")
    except Exception:
        groq_key = settings.get("groq_api_key", "")

    st.markdown("""<div style='padding: 8px 0 24px;'>
        <div style='font-family: JetBrains Mono, monospace; font-size: 11px; color: #3d5a80; letter-spacing: 0.1em;'>AI TOOL</div>
        <h1 style='font-size: 28px; margin: 4px 0 0; color: #e8eaf0;'>Resume Rewriter</h1>
        <p style='color: #5a7090; font-size: 14px; margin: 4px 0 0;'>Upload your old resume (PDF, DOCX, or PPTX) → AI rewrites it → Download a polished DOCX.</p>
    </div>""", unsafe_allow_html=True)

    if not groq_key:
        st.error("⚠️ No Groq API key found. Ask your admin to configure it.")
        return

    col1, col2 = st.columns([1, 1])

    with col1:
        st.markdown('<div class="section-header">STEP 1 — UPLOAD YOUR RESUME</div>', unsafe_allow_html=True)
        uploaded = st.file_uploader(
            "Upload resume",
            type=["pdf", "docx", "doc", "pptx", "ppt"],
            label_visibility="collapsed"
        )

        st.markdown('<div class="section-header">STEP 2 — TARGET JOB (OPTIONAL)</div>', unsafe_allow_html=True)
        jd_text = st.text_area(
            "Paste job description to tailor resume",
            height=180,
            placeholder="Paste the job description here to make the AI tailor your resume specifically to this role. Leave blank for a general rewrite.",
            label_visibility="collapsed"
        )

        st.markdown('<div class="section-header">STEP 3 — REWRITE MODE</div>', unsafe_allow_html=True)
        mode = st.radio(
            "Mode",
            ["🔄 Full Rewrite — rewrite the entire resume",
             "🎯 Targeted — rewrite summary & skills only"],
            label_visibility="collapsed"
        )
        rewrite_mode = "full" if "Full" in mode else "targeted"

        go = st.button("🤖 Rewrite My Resume", disabled=uploaded is None)

    with col2:
        st.markdown('<div class="section-header">PREVIEW & DOWNLOAD</div>', unsafe_allow_html=True)

        if "rewritten_resume" in st.session_state:
            data = st.session_state["rewritten_resume"]

            # Preview
            st.markdown(f"""<div style='background:#111827; border:1px solid #22d3a5; border-radius:8px; padding:20px; font-size:13px; line-height:1.8;'>
                <div style='font-size:18px; font-weight:700; color:#e8eaf0; margin-bottom:8px;'>{data.get("candidate_name","")}</div>
                <div style='color:#5a7090; font-size:11px; font-family: JetBrains Mono; margin-bottom:12px;'>PROFESSIONAL SUMMARY</div>
                <div style='color:#c8d0e0; margin-bottom:16px;'>{data.get("professional_summary","")}</div>
                <div style='color:#5a7090; font-size:11px; font-family: JetBrains Mono; margin-bottom:8px;'>KEY SKILLS</div>
                <div style='color:#22d3a5; font-family: JetBrains Mono; font-size:11px;'>{" · ".join(data.get("skills",[])[:12])}</div>
            </div>""", unsafe_allow_html=True)

            if data.get("keywords_added"):
                st.markdown(f"""<div style='background:#0f3a2a; border:1px solid #22d3a5; border-radius:6px;
                    padding:10px 14px; margin-top:10px; font-family: JetBrains Mono; font-size:11px; color:#22d3a5;'>
                    ✨ NEW KEYWORDS ADDED: {" · ".join(data["keywords_added"][:8])}
                </div>""", unsafe_allow_html=True)

            # Download button
            docx_bytes = build_docx(data)
            fname = f"{data.get('candidate_name','Resume').replace(' ','_')}_AI_Updated.docx"
            st.download_button(
                label="📥 Download Updated Resume (.docx)",
                data=docx_bytes,
                file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )

            exp_count = len(data.get("work_experience", []))
            skills_count = len(data.get("skills", []))
            kw_count = len(data.get("keywords_added", []))
            c1s, c2s, c3s = st.columns(3)
            c1s.metric("Jobs in resume", exp_count)
            c2s.metric("Skills listed", skills_count)
            c3s.metric("New keywords", kw_count)
        else:
            st.markdown("""<div style='background:#111827; border:1px solid #1e2d4a; border-radius:8px;
                padding:40px 20px; text-align:center; color:#3d5a80;'>
                <div style='font-size:40px; margin-bottom:12px;'>📄</div>
                <div style='font-family: JetBrains Mono; font-size:12px;'>Upload a resume and click Rewrite to see the AI-updated version here</div>
            </div>""", unsafe_allow_html=True)

    # ── Run rewrite ──
    if go and uploaded:
        file_bytes = uploaded.read()
        with st.spinner("📖 Reading your resume..."):
            resume_text = extract_text(file_bytes, uploaded.name)

        if not resume_text or len(resume_text) < 50:
            st.error("Could not extract text from this file. Try a different format.")
            return

        with st.spinner("🤖 AI is rewriting your resume... (10-20 seconds)"):
            try:
                result = ai_rewrite_resume(
                    resume_text=resume_text,
                    jd_text=jd_text,
                    name=prefs.get("name", ""),
                    groq_key=groq_key,
                    mode=rewrite_mode,
                )
                st.session_state["rewritten_resume"] = result
                st.rerun()
            except Exception as e:
                st.error(f"AI rewrite failed: {e}")
